import yaml
from pptx import Presentation
from pptx.util import Inches

def create_presentation_from_yaml(yaml_file, output_pptx):
    # Load data from YAML file with UTF-8 encoding
    with open(yaml_file, 'r', encoding='utf-8') as file:
        data = yaml.safe_load(file)

    # Create a presentation object
    presentation = Presentation()

    for slide_data in data:
        slide_layout = presentation.slide_layouts[0] if slide_data['type'] == 'title' else presentation.slide_layouts[1]

        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data['title']
        content.text = slide_data.get('content', '')

    # Save the presentation
    presentation.save(output_pptx)

if __name__ == "__main__":
    yaml_file_path = "presentation_structure.yaml"
    output_pptx_path = "my_presentation.pptx"
    create_presentation_from_yaml(yaml_file_path, output_pptx_path)
